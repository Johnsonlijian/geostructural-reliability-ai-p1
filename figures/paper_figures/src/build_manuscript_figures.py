"""Manuscript figure package, publication style.

Generates deterministic, editable-first manuscript figures. No generative image model, web icon,
or stock artwork is used.

Style contract (R15 figure-style correction):
- Fig. 1 is the single schematic/mechanism figure (flat scientific schematic).
- Figs. 2-5 are journal-style data figures: white background, thin axes, no
  in-figure title banners, no rounded-card panels, no drop shadows, no slogans.
- All figures are designed at 180 mm print width with 7-8 pt text so that
  \\includegraphics[width=\\linewidth] reproduces near-nominal font sizes.
"""

from __future__ import annotations

import csv
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import yaml
from matplotlib.lines import Line2D
from matplotlib.patches import Circle, FancyArrowPatch, Polygon, Rectangle

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    HAVE_PPTX = True
except ImportError:  # pptx previews are optional; SVG/PDF are canonical
    HAVE_PPTX = False


SCRIPT = Path(__file__).resolve()
FIG_ROOT = SCRIPT.parents[1]
PROJECT = SCRIPT.parents[3]
OUT = FIG_ROOT / "output"
DATA = FIG_ROOT / "data"
BLENDER = FIG_ROOT / "blender"

# Okabe-Ito colorblind-safe palette
BLUE = "#0072B2"     # SPT / primary dataset
PURPLE = "#CC79A7"   # CPT / secondary dataset
GREEN = "#009E73"    # stable / gain excluded
AMBER = "#E69F00"    # critical band / caveat
RED = "#D55E00"      # failed / single-event caveat
GRAY = "#7f7f7f"
DARK = "#1f2933"
SOIL = "#d6bd8f"
SAND = "#ead7a0"
WATER = "#5dade2"

FULL_W = 7.05  # inches, ~180 mm double-column print width

KEY = {
    "random_optimism": {"SPT": 0.047, "CPT": 0.117},
    "cpt2024": {
        "inverse": {"d_auc": -0.071, "auc_lo": -0.141, "auc_hi": -0.009, "d_nll": -0.092, "nll_lo": -0.130, "nll_hi": -0.051},
        "measured": {"d_auc": -0.051, "auc_lo": -0.084, "auc_hi": -0.009, "d_nll": -0.116, "nll_lo": -0.163, "nll_hi": -0.055},
    },
    "spt_caveats": {
        "source": {"d_auc": 0.041, "lo": 0.004, "hi": 0.069},
        "source_event": {"d_auc": 0.035, "lo": -0.009, "hi": 0.056},
    },
    "conformal": {
        "SPT": {"coverage": 0.925, "set": 1.535, "two_label": 0.542, "critical_two": 0.696},
        "CPT": {"coverage": 0.897, "set": 1.691, "two_label": 0.693, "critical_two": 0.807},
    },
    "alpha": {
        "SPT": [(0.05, 0.942, 1.664, 0.674, 0.833), (0.10, 0.925, 1.535, 0.542, 0.696), (0.20, 0.838, 1.320, 0.341, 0.447)],
        "CPT": [(0.05, 0.932, 1.806, 0.809, 0.871), (0.10, 0.897, 1.691, 0.693, 0.807), (0.20, 0.801, 1.418, 0.432, 0.617)],
    },
    "groundwater": {"label": ["surface", "actual", "deep/dry"], "fs": [0.64, 0.90, 1.25], "risk": [0.73, 0.57, 0.30]},
    # bin-level margin residuals mirrored from the archived analysis output
    # code/data/processed/groundwater_residual_stratification.json (means and 95% cluster-bootstrap CIs)
    "groundwater_residuals": {
        "bins": ["<=2 m", "2-5 m", "5-10 m"],
        "SPT": {"mean": [-0.00246, -0.01803, 0.05729], "lo": [-0.12213, -0.16396, -0.31517], "hi": [0.08725, 0.12135, 0.43102], "n": [122, 81, 5]},
        "CPT": {"mean": [-0.03632, 0.03756, 0.08149], "lo": [-0.16153, -0.12534, -0.14409], "hi": [0.08780, 0.21432, 0.39698], "n": [131, 130, 14]},
    },
    "slope": {
        "auc": 0.895,
        "dummy_delta_auc": -0.012,
        "mid_below_anchor": 0.953,
        "mid_q": [0.763, 0.879, 0.998],
        "points": [
            ("CB1 fail", 0.908, 1, "CB1"),
            ("CB1 stable", 0.639, 0, "CB1"),
            ("Flume 1 fail", 0.992, 1, "USGS flume"),
            ("Flume 2 fail", 1.047, 1, "USGS flume"),
            ("Flume 3 control", 0.852, 1, "domain control"),
            ("BALT1 stable", 0.805, 0, "BALT"),
            ("BALT2 stable", 0.170, 0, "BALT"),
            ("BALT3 stable", 0.799, 0, "BALT"),
        ],
    },
}


def setup() -> None:
    for sub in ["svg", "pdf", "png", "pptx", "blend"]:
        (OUT / sub).mkdir(parents=True, exist_ok=True)
    DATA.mkdir(parents=True, exist_ok=True)
    BLENDER.mkdir(parents=True, exist_ok=True)
    plt.rcParams.update(
        {
            "font.family": "sans-serif",
            "font.sans-serif": ["Arial", "Helvetica", "DejaVu Sans"],
            "font.size": 7.5,
            "axes.labelsize": 8.0,
            "axes.titlesize": 8.0,
            "xtick.labelsize": 7.0,
            "ytick.labelsize": 7.0,
            "legend.fontsize": 6.8,
            "legend.frameon": False,
            "axes.linewidth": 0.6,
            "xtick.major.width": 0.6,
            "ytick.major.width": 0.6,
            "xtick.major.size": 2.5,
            "ytick.major.size": 2.5,
            "xtick.direction": "out",
            "ytick.direction": "out",
            "lines.linewidth": 1.0,
            "lines.markersize": 3.5,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "figure.dpi": 160,
            "savefig.dpi": 600,
            "savefig.facecolor": "white",
            "svg.fonttype": "none",
            "pdf.fonttype": 42,
            "ps.fonttype": 42,
        }
    )


def save(fig: plt.Figure, stem: str, title: str) -> None:
    for ext in ["svg", "pdf", "png"]:
        fig.savefig(OUT / ext / f"{stem}.{ext}", bbox_inches="tight")
    make_pptx_preview(stem, title)
    plt.close(fig)


def make_pptx_preview(stem: str, title: str) -> None:
    if not HAVE_PPTX:
        return
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.12), Inches(12.8), Inches(0.35))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.name = "Arial"
    png = OUT / "png" / f"{stem}.png"
    slide.shapes.add_picture(str(png), Inches(0.20), Inches(0.55), width=Inches(12.9))
    note = slide.shapes.add_textbox(Inches(0.25), Inches(7.05), Inches(12.6), Inches(0.25))
    note.text_frame.text = "Editable scientific source is the SVG/PDF plus archived Python script and source data. PNG is a preview layer."
    note.text_frame.paragraphs[0].font.size = Pt(8)
    prs.save(OUT / "pptx" / f"{stem}_preview.pptx")


def panel_letter(ax, letter: str, dx_pt: float = -26.0, dy_pt: float = 2.0) -> None:
    """Bold lowercase panel letter, offset in points from the axes top-left corner."""
    ax.annotate(
        letter,
        xy=(0, 1),
        xycoords="axes fraction",
        xytext=(dx_pt, dy_pt),
        textcoords="offset points",
        fontsize=9.5,
        fontweight="bold",
        ha="left",
        va="bottom",
    )


def arrow(ax, xy1, xy2, color=DARK, lw=0.9, style="->", mutation=8, alpha=1.0):
    ax.add_patch(FancyArrowPatch(xy1, xy2, arrowstyle=style, mutation_scale=mutation, lw=lw, color=color, alpha=alpha))


def draw_soil_cutaway(ax, slope=False):
    """Flat scientific schematic; used in Fig. 1a (full) and Fig. 5a (slope)."""
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    if slope:
        ground = Polygon([(0.08, 0.25), (0.86, 0.25), (0.86, 0.70), (0.08, 0.92)], fc=SOIL, ec=DARK, lw=0.7)
        ax.add_patch(ground)
        ax.plot([0.15, 0.82], [0.43, 0.57], color=RED, lw=1.4)
        ax.text(0.55, 0.62, "slip surface", fontsize=6.5, color=RED)
        for x in np.linspace(0.18, 0.76, 7):
            arrow(ax, (x, 0.99), (x - 0.02, 0.84), color=WATER, lw=0.7, mutation=6)
        ax.text(0.13, 0.95, "rainfall", color=BLUE, fontsize=6.8)
        for x in [0.32, 0.48, 0.64]:
            arrow(ax, (x, 0.78), (x + 0.02, 0.56), color=BLUE, lw=0.7, mutation=6)
        ax.text(0.10, 0.16, "effective-stress capacity\ncoordinate $X+Y$", fontsize=6.8, color=DARK)
    else:
        ax.add_patch(Polygon([(0.10, 0.15), (0.82, 0.15), (0.92, 0.72), (0.22, 0.72)], fc=SOIL, ec=DARK, lw=0.7))
        ax.add_patch(Polygon([(0.14, 0.30), (0.86, 0.30), (0.82, 0.45), (0.18, 0.45)], fc=SAND, ec="#b7945d", lw=0.5))
        ax.add_patch(Polygon([(0.16, 0.53), (0.88, 0.53), (0.86, 0.60), (0.18, 0.60)], fc=WATER, ec=BLUE, lw=0.5, alpha=0.35))
        ax.text(0.76, 0.62, "water table", color=BLUE, fontsize=6.5)
        ax.plot([0.12, 0.90], [0.78, 0.78], color="#4d5562", lw=0.7)
        ax.add_patch(Rectangle((0.30, 0.78), 0.10, 0.10, fc="#dde4ec", ec=DARK, lw=0.6))
        ax.text(0.25, 0.905, "bridge / lifeline", fontsize=6.5, color=DARK)
        for y in [0.075, 0.035]:
            ax.plot(np.linspace(0.12, 0.82, 80), y + 0.012 * np.sin(np.linspace(0, 18, 80)), color=GRAY, lw=0.7)
        ax.text(0.12, 0.00, "seismic demand", color=GRAY, fontsize=6.5)
        ax.add_patch(Polygon([(0.58, 0.83), (0.61, 0.83), (0.61, 0.38), (0.595, 0.32), (0.58, 0.38)], fc="#e8eef6", ec=DARK, lw=0.7))
        ax.text(0.625, 0.73, "CPT", fontsize=6.5, color=DARK)
        ax.plot([0.43, 0.43], [0.82, 0.42], color=DARK, lw=1.2)
        ax.add_patch(Rectangle((0.39, 0.82), 0.08, 0.035, fc="#d9dee7", ec=DARK, lw=0.6))
        ax.text(0.31, 0.83, "SPT", fontsize=6.5, color=DARK)
        for x, scale in [(0.24, 0.22), (0.33, 0.16), (0.50, 0.10)]:
            arrow(ax, (x, 0.66), (x, 0.66 - scale), color=BLUE, lw=0.9, mutation=7)
        ax.text(0.14, 0.70, "decreasing $\\sigma'_v$", color=BLUE, fontsize=6.5)
        arrow(ax, (0.20, 0.38), (0.47, 0.38), color=RED, lw=1.2, mutation=9)
        arrow(ax, (0.79, 0.38), (0.57, 0.38), color=GREEN, lw=1.2, mutation=9)
        ax.text(0.20, 0.415, "CSR demand", color=RED, fontsize=6.5)
        ax.text(0.63, 0.415, "CRR resistance", color=GREEN, fontsize=6.5)
        ax.add_patch(Circle((0.52, 0.38), 0.035, fc="white", ec=DARK, lw=0.8))
        ax.text(0.52, 0.38, "FS", ha="center", va="center", fontsize=6.8, color=DARK)
        ax.text(0.30, 0.22, "non-refitted margin  $s=\\ln(\\mathrm{CSR}/\\mathrm{CRR})$", fontsize=6.8, color=DARK)


def fig1():
    """Fig. 1: the single permitted mechanism/protocol schematic (flat, no cards)."""
    fig = plt.figure(figsize=(FULL_W, 3.0))
    gs = fig.add_gridspec(1, 2, width_ratios=[1.45, 1.0], wspace=0.06)

    ax = fig.add_subplot(gs[0, 0])
    draw_soil_cutaway(ax)
    panel_letter(ax, "a", dx_pt=-2)

    axr = fig.add_subplot(gs[0, 1])
    axr.set_xlim(0, 1)
    axr.set_ylim(0, 1)
    axr.axis("off")
    panel_letter(axr, "b", dx_pt=-2)
    steps = [
        ("non-refitted physics margin", None),
        ("grouped transfer\n(event / source / site)", None),
        ("paired gain-exclusion audit", "gain-exclusion record (Fig. 2)"),
        ("residual caveat retention", "transfer-axis caveat map (Fig. 2d)"),
        ("mechanism-band\nconformal review", "two-label review sets (Fig. 3)"),
    ]
    box_w, box_h, gap = 0.52, 0.115, 0.058
    y0 = 0.97
    for i, (txt, out_note) in enumerate(steps):
        y = y0 - i * (box_h + gap) - box_h
        axr.add_patch(Rectangle((0.02, y), box_w, box_h, fc="white", ec=DARK, lw=0.8))
        axr.text(0.02 + box_w / 2, y + box_h / 2, txt, ha="center", va="center", fontsize=6.6, color=DARK)
        if i < len(steps) - 1:
            arrow(axr, (0.02 + box_w / 2, y - 0.004), (0.02 + box_w / 2, y - gap + 0.004), color=DARK, lw=0.8, mutation=7)
        if out_note:
            arrow(axr, (0.02 + box_w + 0.005, y + box_h / 2), (0.02 + box_w + 0.05, y + box_h / 2), color=GRAY, lw=0.7, mutation=6)
            axr.text(0.02 + box_w + 0.06, y + box_h / 2, out_note, ha="left", va="center", fontsize=6.2, color=GRAY)
    save(fig, "Fig1_hero_mechanism", "Fig. 1 Mechanism-anchored validation protocol")


def draw_split_schematic(ax):
    """Fig. 2a: clustered records and the two split designs (no card backgrounds)."""
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    rng = np.random.default_rng(7)
    centers = [(0.16, 0.62), (0.30, 0.36), (0.40, 0.66)]
    cluster_pts = []
    for cx, cy in centers:
        pts = np.column_stack([cx + 0.055 * rng.standard_normal(9), cy + 0.060 * rng.standard_normal(9)])
        cluster_pts.append(pts)

    # left scene: random split (test points drawn from every cluster)
    for k, pts in enumerate(cluster_pts):
        ax.scatter(pts[:, 0], pts[:, 1], s=7, color=GRAY, alpha=0.8, linewidths=0)
        test = pts[rng.choice(len(pts), 3, replace=False)]
        ax.scatter(test[:, 0], test[:, 1], s=12, facecolor="none", edgecolor=RED, linewidths=0.8)
    ax.text(0.26, 0.10, "random split\ntest cases mix clusters", ha="center", fontsize=6.6, color=DARK)

    # right scene: same clusters shifted right, grouped split holds out one block
    dx = 0.46
    for k, pts in enumerate(cluster_pts):
        p = pts + np.array([dx, 0.0])
        if k == 2:
            ax.scatter(p[:, 0], p[:, 1], s=12, facecolor="none", edgecolor=RED, linewidths=0.8)
        else:
            ax.scatter(p[:, 0], p[:, 1], s=7, color=GRAY, alpha=0.8, linewidths=0)
    held = cluster_pts[2] + np.array([dx, 0.0])
    x0, y0 = held[:, 0].min() - 0.035, held[:, 1].min() - 0.045
    ax.add_patch(Rectangle((x0, y0), held[:, 0].max() - x0 + 0.035, held[:, 1].max() - y0 + 0.045,
                           fc="none", ec=DARK, lw=0.8))
    ax.text(0.72, 0.10, "grouped split\nholds out whole blocks", ha="center", fontsize=6.6, color=DARK)
    handles = [
        Line2D([0], [0], marker="o", color="none", markerfacecolor=GRAY, markersize=3, label="training case"),
        Line2D([0], [0], marker="o", color="none", markerfacecolor="none", markeredgecolor=RED, markersize=3.6, label="test case"),
    ]
    ax.legend(handles=handles, loc="upper center", bbox_to_anchor=(0.5, 1.06), ncol=2,
              fontsize=6.2, handletextpad=0.3, columnspacing=0.8)


def fig2():
    fig, axs = plt.subplots(2, 2, figsize=(FULL_W, 5.35),
                            gridspec_kw={"width_ratios": [1.0, 1.08]})
    fig.subplots_adjust(hspace=0.44, wspace=0.54)

    ax = axs[0, 0]
    draw_split_schematic(ax)
    panel_letter(ax, "a", dx_pt=-6)

    ax = axs[0, 1]
    panel_letter(ax, "b")
    datasets = ["SPT", "CPT"]
    grouped = [0.865, 0.613]
    random = [0.912, 0.730]
    margin = [0.923, 0.750]
    for i in range(2):
        col = PURPLE if i else BLUE
        ax.plot([0, 1], [grouped[i], random[i]], marker="o", color=col, lw=1.0)
        ax.text(1.05, random[i], f"{datasets[i]} +{random[i]-grouped[i]:.3f}", va="center", fontsize=6.6, color=col)
    ax.scatter([0, 0], margin, marker="s", s=14, color=GREEN, zorder=3, label="mechanism margin")
    ax.set_xlim(-0.15, 1.55)
    ax.set_xticks([0, 1])
    ax.set_xticklabels(["grouped ML", "random ML"])
    ax.set_ylabel("AUC")
    ax.set_ylim(0.55, 0.98)
    ax.legend(loc="lower left", handletextpad=0.3)

    ax = axs[1, 0]
    panel_letter(ax, "c", dx_pt=-44)
    models = ["logistic", "random forest", "gradient boost", "hist. GBT", "neural net"]
    spt = [-0.025, -0.044, -0.140, -0.047, -0.758]
    cpt = [-0.133, -0.101, -0.254, -0.465, -1.649]
    y = np.arange(len(models))
    ax.axvline(0, color=DARK, lw=0.6)
    ax.scatter(spt, y + 0.12, s=12, color=BLUE, label="SPT")
    ax.scatter(cpt, y - 0.12, s=12, color=PURPLE, label="CPT")
    ax.set_yticks(y)
    ax.set_yticklabels(models)
    ax.set_xlabel("likelihood gain vs margin (positive favors candidate)")
    ax.annotate("favors margin", xy=(-1.50, 1.5), xytext=(-0.92, 1.5),
                arrowprops=dict(arrowstyle="->", lw=0.7, color=GRAY),
                fontsize=6.6, color=GRAY, va="center")
    ax.legend(loc="lower left", handletextpad=0.3)

    ax = axs[1, 1]
    panel_letter(ax, "d", dx_pt=-42)
    rows = [
        ("CPT-2024 inverse", -0.071, -0.141, -0.009, "gain excluded", GREEN),
        ("CPT-2024 measured", -0.051, -0.084, -0.009, "gain excluded", GREEN),
        ("CPT-2021 event", -0.051, -0.128, -0.011, "gain excluded", GREEN),
        ("Vs-2013 event", -0.029, -0.061, -0.002, "gain excluded", GREEN),
        ("SPT source", 0.041, 0.004, 0.069, "caveat retained", AMBER),
        ("SPT source-event", 0.035, -0.009, 0.056, "caveat retained", AMBER),
        ("Nisqually inverse", 0.030, 0.000, 0.130, "single-event\ncaveat", RED),
    ]
    y = np.arange(len(rows))[::-1]
    for yy, (lab, val, lo, hi, status, col) in zip(y, rows):
        ax.plot([lo, hi], [yy, yy], color=col, lw=1.1, solid_capstyle="butt")
        ax.scatter([val], [yy], color=col, s=14, zorder=3)
        ax.text(0.225, yy, status, va="center", ha="right", fontsize=6.2, color=col)
    ax.axvline(0.02, color=GRAY, ls="--", lw=0.7)
    ax.set_ylim(-0.5, 7.1)
    ax.text(0.024, 6.95, "+0.02 practical\ngain threshold", fontsize=6.0, color=GRAY, va="top")
    ax.axvline(0, color=DARK, lw=0.6)
    ax.set_yticks(y)
    ax.set_yticklabels([
        "CPT-24 inverse",
        "CPT-24 measured",
        "CPT-21 event",
        "Vs-13 event",
        "SPT source",
        "SPT source-event",
        "Nisqually inv.",
    ])
    ax.tick_params(axis="y", labelsize=6.8, pad=1.5)
    ax.set_xlabel(r"$\Delta$AUC, non-margin learner minus margin")
    ax.set_xlim(-0.16, 0.23)
    save(fig, "Fig2_transfer_evidence_map", "Fig. 2 Validation stress-test evidence")


def fig3():
    fig, axs = plt.subplots(2, 2, figsize=(FULL_W, 5.1))
    fig.subplots_adjust(hspace=0.42, wspace=0.32)

    ax = axs[0, 0]
    panel_letter(ax, "a", dx_pt=-8)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axvspan(0.3, 0.7, color=AMBER, alpha=0.15, lw=0)
    for xb in (0.3, 0.7):
        ax.axvline(xb, color=GRAY, ls="--", lw=0.7)
    ax.text(0.15, 0.55, r"$\{0\}$", ha="center", fontsize=8.0, color=GREEN)
    ax.text(0.50, 0.55, r"$\{0,1\}$", ha="center", fontsize=8.0, color=AMBER)
    ax.text(0.85, 0.55, r"$\{1\}$", ha="center", fontsize=8.0, color=RED)
    ax.text(0.15, 0.40, "singleton\ndecision", ha="center", va="top", fontsize=6.4, color=GRAY)
    ax.text(0.50, 0.40, "two-label\nreview set", ha="center", va="top", fontsize=6.4, color=GRAY)
    ax.text(0.85, 0.40, "singleton\ndecision", ha="center", va="top", fontsize=6.4, color=GRAY)
    ax.text(0.50, 0.92, "critical band [0.3, 0.7]", ha="center", fontsize=6.6, color=DARK)
    ax.set_xlabel("calibrated mechanism probability")
    ax.set_yticks([])
    ax.spines["left"].set_visible(False)

    ax = axs[0, 1]
    panel_letter(ax, "b")
    for name, col in [("SPT", BLUE), ("CPT", PURPLE)]:
        arr = KEY["alpha"][name]
        set_size = [r[2] for r in arr]
        cov = [r[1] for r in arr]
        ax.plot(set_size, cov, "-o", color=col, label=name, lw=1.0)
        for a, c, s, _, _ in arr:
            ax.text(s + 0.012, c, rf"$\alpha$={a:.2f}", fontsize=6.0, color=col, va="center")
    ax.axhline(0.90, color=GRAY, ls="--", lw=0.7)
    ax.text(1.86, 0.903, "0.90", fontsize=6.0, color=GRAY, ha="right")
    ax.set_xlabel("mean prediction-set size")
    ax.set_ylabel("coverage")
    ax.set_xlim(1.25, 1.95)
    ax.set_ylim(0.78, 0.96)
    ax.legend(loc="lower right", handletextpad=0.3)

    ax = axs[1, 0]
    panel_letter(ax, "c")
    metrics = ["coverage", "two-label", "critical\ntwo-label", "event\nundercov."]
    spt = [0.922, 0.515, 0.790, 0.286]
    cpt = [0.913, 0.728, 0.879, 0.250]
    xx = np.arange(len(metrics))
    ax.bar(xx - 0.18, spt, 0.36, color=BLUE, label="SPT")
    ax.bar(xx + 0.18, cpt, 0.36, color=PURPLE, label="CPT")
    ax.set_xticks(xx)
    ax.set_xticklabels(metrics)
    ax.set_ylim(0, 1.0)
    ax.set_ylabel("median over 100 grouped splits")
    ax.legend(loc="upper right", handletextpad=0.3)

    ax = axs[1, 1]
    panel_letter(ax, "d")
    review_fn5 = np.array([0.25, 0.50])
    spt_fn5 = np.array([0.469, 0.338])
    cpt_fn5 = np.array([0.699, 0.512])
    ax.plot(review_fn5, spt_fn5, "-o", color=BLUE, label=r"SPT, $c_{\mathrm{fn}}$=5", lw=1.0)
    ax.plot(review_fn5, cpt_fn5, "-o", color=PURPLE, label=r"CPT, $c_{\mathrm{fn}}$=5", lw=1.0)
    ax.scatter([0.50], [0.800], marker="D", s=18, color=BLUE, label=r"SPT, $c_{\mathrm{fn}}$=10")
    ax.scatter([0.50], [1.244], marker="D", s=18, color=PURPLE, label=r"CPT, $c_{\mathrm{fn}}$=10")
    ax.set_xlabel("review cost")
    ax.set_ylabel("expected cost reduction")
    ax.set_xticks([0.25, 0.50])
    ax.legend(loc="upper left", ncol=2, columnspacing=0.8, handletextpad=0.3)
    save(fig, "Fig3_conformal_cockpit", "Fig. 3 Conformal review reliability")


def fig4():
    fig = plt.figure(figsize=(FULL_W, 2.9))
    gs = fig.add_gridspec(1, 2, width_ratios=[1.2, 1.0], wspace=0.22)

    ax = fig.add_subplot(gs[0, 0])
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    panel_letter(ax, "a", dx_pt=-2)
    xs = [0.06, 0.37, 0.68]
    water_y = [0.74, 0.57, 0.34]
    for i, x0 in enumerate(xs):
        ax.add_patch(Rectangle((x0, 0.20), 0.24, 0.66, fc=SOIL, ec=DARK, lw=0.7))
        ax.add_patch(Rectangle((x0, 0.20), 0.24, water_y[i] - 0.20, fc=WATER, ec="none", alpha=0.30))
        ax.plot([x0, x0 + 0.24], [water_y[i], water_y[i]], color=BLUE, lw=1.2)
        ax.text(x0 + 0.12, 0.90, KEY["groundwater"]["label"][i], ha="center", fontsize=7.0, color=DARK)
        ax.text(x0 + 0.12, 0.105, f"FS {KEY['groundwater']['fs'][i]:.2f}", ha="center", fontsize=7.0, color=DARK)
        ax.text(x0 + 0.12, 0.025, f"Pr(FS<1) {KEY['groundwater']['risk'][i]:.2f}", ha="center", fontsize=6.5, color=GRAY)
    arrow(ax, (0.345, water_y[0]), (0.365, water_y[1] + 0.012), color=BLUE, lw=0.8, mutation=7)
    arrow(ax, (0.655, water_y[1]), (0.675, water_y[2] + 0.012), color=BLUE, lw=0.8, mutation=7)
    ax.text(0.06, 0.955, "water-table counterfactual", fontsize=6.8, color=GRAY)

    ax = fig.add_subplot(gs[0, 1])
    panel_letter(ax, "b")
    gw = KEY["groundwater_residuals"]
    x = np.arange(3)
    for name, col, marker, off in [("SPT", BLUE, "o", -0.09), ("CPT", PURPLE, "s", 0.09)]:
        d = gw[name]
        mean = np.array(d["mean"])
        lo = np.array(d["lo"])
        hi = np.array(d["hi"])
        ax.errorbar(x + off, mean, yerr=[mean - lo, hi - mean], fmt=marker, ms=3.2, color=col,
                    capsize=2, lw=0.9, label=name)
    ax.axhline(0, color=DARK, lw=0.6)
    ax.set_xticks(x)
    ax.set_xticklabels([r"$\leq$2", "2-5", "5-10"])
    for i in range(3):
        ax.text(i, -0.46, f"n={gw['SPT']['n'][i]}/{gw['CPT']['n'][i]}", ha="center", fontsize=6.0, color=GRAY)
    ax.set_ylim(-0.52, 0.52)
    ax.set_xlabel("ground-water depth bin (m)")
    ax.set_ylabel("mean margin residual (95% CI)")
    ax.legend(loc="upper left", handletextpad=0.3)
    save(fig, "Fig4_groundwater_cutaway", "Fig. 4 Groundwater mechanism and residual boundary")


def fig5():
    fig = plt.figure(figsize=(FULL_W, 2.7))
    gs = fig.add_gridspec(1, 3, width_ratios=[0.95, 1.05, 0.95], wspace=0.38)

    ax = fig.add_subplot(gs[0, 0])
    draw_soil_cutaway(ax, slope=True)
    panel_letter(ax, "a", dx_pt=-2)

    ax = fig.add_subplot(gs[0, 1])
    panel_letter(ax, "b", dx_pt=-50)
    points = KEY["slope"]["points"]
    ymap = {"CB1": 3, "USGS flume": 2, "BALT": 1, "domain control": 0}
    marker_map = {"CB1": "o", "USGS flume": "o", "BALT": "o", "domain control": "^"}
    jitter_map = {
        "CB1 fail": 0.10,
        "CB1 stable": -0.10,
        "Flume 1 fail": 0.12,
        "Flume 2 fail": -0.12,
        "Flume 3 control": 0.00,
        "BALT1 stable": -0.10,
        "BALT2 stable": 0.00,
        "BALT3 stable": 0.10,
    }
    q05, q50, q95 = KEY["slope"]["mid_q"]
    ax.axvspan(0.805, 0.908, color=AMBER, alpha=0.15, lw=0)
    for name, xy, label, group in points:
        col = RED if label == 1 else GREEN
        ax.scatter(xy, ymap[group] + jitter_map[name], s=16, color=col, edgecolor=DARK,
                   lw=0.3, marker=marker_map[group], zorder=3)
    ax.axvline(1.0, color=DARK, lw=0.8)
    ax.set_yticks([3, 2, 1, 0])
    ax.set_yticklabels(["CB1", "USGS flume", "BALT", "domain control"])
    ax.set_xlabel("slope capacity coordinate $X+Y$")
    ax.set_xlim(0.10, 1.12)
    ax.set_ylim(-1.45, 3.55)
    handles = [
        Line2D([0], [0], marker="o", color="none", markerfacecolor=GREEN, markeredgecolor=DARK, markersize=3.6, label="stable"),
        Line2D([0], [0], marker="o", color="none", markerfacecolor=RED, markeredgecolor=DARK, markersize=3.6, label="failed"),
        Line2D([0], [0], marker="^", color="none", markerfacecolor=RED, markeredgecolor=DARK, markersize=3.6, label="domain control"),
        Rectangle((0, 0), 1, 1, fc=AMBER, ec="none", alpha=0.15, label="transition band"),
    ]
    ax.legend(handles=handles, fontsize=6.0, loc="lower left", ncol=2, columnspacing=0.7,
              handletextpad=0.3, borderaxespad=0.2)

    ax = fig.add_subplot(gs[0, 2])
    panel_letter(ax, "c")
    ax.axvspan(q05, q95, color=AMBER, alpha=0.20, lw=0)
    ax.axvline(q50, color=AMBER, lw=1.4)
    ax.axvline(1.0, color=DARK, lw=0.8)
    ax.text(q50, 1.02, f"median {q50:.3f}", ha="center", va="bottom", fontsize=6.5, color=DARK)
    ax.text(1.0, 1.10, "$X+Y=1$ anchor", ha="center", va="bottom", fontsize=6.5, color=DARK)
    ax.text(0.70, 0.50, f"P(midpoint < anchor)\n= {KEY['slope']['mid_below_anchor']:.3f}", fontsize=7.0, color=DARK)
    ax.text(0.70, 0.18, "5-95% band", fontsize=6.5, color=GRAY)
    ax.set_xlim(0.68, 1.10)
    ax.set_ylim(0, 1)
    ax.set_yticks([])
    ax.spines["left"].set_visible(False)
    ax.set_xlabel("transition midpoint")
    save(fig, "Fig5_cross_hazard_vignette", "Fig. 5 Cross-hazard extension vignette")


def write_source_data() -> None:
    (DATA / "fig1_key_values.yml").write_text(yaml.safe_dump(KEY, sort_keys=False), encoding="utf-8")
    with (DATA / "fig2_transfer_stress_tests.csv").open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["panel", "item", "dataset", "metric", "value", "lo", "hi", "status_or_note"])
        for dataset, grouped, random, margin in [
            ("SPT", 0.865, 0.912, 0.923),
            ("CPT", 0.613, 0.730, 0.750),
        ]:
            w.writerow(["b", "grouped ML", dataset, "AUC", grouped, "", "", "grouped validation"])
            w.writerow(["b", "random ML", dataset, "AUC", random, "", "", "100-repeat median"])
            w.writerow(["b", "mechanism", dataset, "AUC", margin, "", "", "margin baseline"])
        for model, spt_gain, cpt_gain in [
            ("logistic", -0.025, -0.133),
            ("random forest", -0.044, -0.101),
            ("gradient boost", -0.140, -0.254),
            ("hist. GBT", -0.047, -0.465),
            ("neural network", -0.758, -1.649),
        ]:
            w.writerow(["c", model, "SPT", "NLL gain vs margin", spt_gain, "", "", "positive would favor candidate"])
            w.writerow(["c", model, "CPT", "NLL gain vs margin", cpt_gain, "", "", "positive would favor candidate"])
        for row in [
            ("CPT-2024 inverse", -0.071, -0.141, -0.009, "gain excluded"),
            ("CPT-2024 measured", -0.051, -0.084, -0.009, "gain excluded"),
            ("CPT-2021 event", -0.051, -0.128, -0.011, "gain excluded"),
            ("Vs-2013 event", -0.029, -0.061, -0.002, "gain excluded; same-population third instrument"),
            ("SPT source", 0.041, 0.004, 0.069, "caveat retained"),
            ("SPT source-event", 0.035, -0.009, 0.056, "caveat retained"),
            ("Nisqually inverse", 0.030, 0.000, 0.130, "single-event caveat"),
        ]:
            item, value, lo, hi, status = row
            w.writerow(["d", item, "", "Delta AUC, non-margin learner minus margin", value, lo, hi, status])
    with (DATA / "fig3_conformal_metrics.csv").open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["panel", "dataset", "metric", "setting", "value", "note"])
        for dataset, arr in KEY["alpha"].items():
            for row in arr:
                alpha, coverage, mean_set_size, two_label, critical_two_label = row
                w.writerow(["b", dataset, "coverage", f"alpha={alpha}", coverage, ""])
                w.writerow(["b", dataset, "mean_set_size", f"alpha={alpha}", mean_set_size, ""])
                w.writerow(["b", dataset, "two_label", f"alpha={alpha}", two_label, ""])
                w.writerow(["b", dataset, "critical_two_label", f"alpha={alpha}", critical_two_label, ""])
        for dataset, values in [
            ("SPT", [("coverage", 0.922), ("two_label", 0.515), ("critical_two_label", 0.790), ("event_undercoverage", 0.286)]),
            ("CPT", [("coverage", 0.913), ("two_label", 0.728), ("critical_two_label", 0.879), ("event_undercoverage", 0.250)]),
        ]:
            for metric, value in values:
                w.writerow(["c", dataset, metric, "100 grouped split median", value, "S18 split sensitivity"])
        for dataset, fn_cost, review_cost, reduction in [
            ("SPT", 5, 0.25, 0.469),
            ("SPT", 5, 0.50, 0.338),
            ("CPT", 5, 0.25, 0.699),
            ("CPT", 5, 0.50, 0.512),
            ("SPT", 10, 0.50, 0.800),
            ("CPT", 10, 0.50, 1.244),
        ]:
            w.writerow(["d", dataset, "expected_cost_reduction", f"false_negative_cost={fn_cost};review_cost={review_cost}", reduction, "S17 exact point"])
        w.writerow(["a", "", "schematic", "critical band [0.3,0.7]", "", "author-defined schematic, not a data panel"])
    with (DATA / "fig4_groundwater_residuals.csv").open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["panel", "dataset", "scenario_or_bin", "metric", "value", "lo", "hi", "n", "note"])
        for row in zip(KEY["groundwater"]["label"], KEY["groundwater"]["fs"], KEY["groundwater"]["risk"]):
            scenario, fs, risk = row
            w.writerow(["a", "", scenario, "median_FS", fs, "", "", "", "water-table counterfactual"])
            w.writerow(["a", "", scenario, "Pr_FS_lt_1", risk, "", "", "", "water-table counterfactual"])
        gw = KEY["groundwater_residuals"]
        for dataset in ["SPT", "CPT"]:
            d = gw[dataset]
            for bin_name, value, lo, hi, n in zip(gw["bins"], d["mean"], d["lo"], d["hi"], d["n"]):
                w.writerow(["b", dataset, bin_name, "mean_margin_residual", value, lo, hi, n,
                            "95% cluster-bootstrap CI; archived in code/data/processed/groundwater_residual_stratification.json"])
    with (DATA / "fig5_slope_vignette.csv").open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["panel", "case_or_metric", "X_plus_Y", "label", "source_group", "value", "note"])
        for row in KEY["slope"]["points"]:
            case, xy, label, source_group = row
            w.writerow(["b", case, xy, label, source_group, "", "verified slope-vignette state"])
        for metric, value in [
            ("midpoint_q05", KEY["slope"]["mid_q"][0]),
            ("midpoint_q50", KEY["slope"]["mid_q"][1]),
            ("midpoint_q95", KEY["slope"]["mid_q"][2]),
            ("P_midpoint_below_anchor", KEY["slope"]["mid_below_anchor"]),
            ("cross_hazard_AUC", KEY["slope"]["auc"]),
            ("hazard_dummy_Delta_AUC", KEY["slope"]["dummy_delta_auc"]),
        ]:
            w.writerow(["c", metric, "", "", "", value, "small-N diagnostic reported in SI"])


def write_provenance() -> None:
    provenance = {}
    for stem, title, scripts, data_files in [
        ("Fig1_hero_mechanism", "Mechanism-anchored validation protocol (single schematic figure)", ["src/build_manuscript_figures.py"], ["data/fig1_key_values.yml"]),
        ("Fig2_transfer_evidence_map", "Validation stress-test evidence", ["src/build_manuscript_figures.py"], ["data/fig2_transfer_stress_tests.csv"]),
        ("Fig3_conformal_cockpit", "Conformal review reliability", ["src/build_manuscript_figures.py"], ["data/fig3_conformal_metrics.csv"]),
        ("Fig4_groundwater_cutaway", "Groundwater mechanism and residual boundary", ["src/build_manuscript_figures.py"], ["data/fig4_groundwater_residuals.csv"]),
        ("Fig5_cross_hazard_vignette", "Cross-hazard extension vignette", ["src/build_manuscript_figures.py"], ["data/fig5_slope_vignette.csv"]),
        ("Fig6_critical_band_closure", "R16/R20 candidate site-paired Vs critical-band closure", ["src/build_manuscript_figures.py"], ["data/fig6_case_level.csv"]),
    ]:
        provenance[stem] = {
            "title": title,
            "data_sources": data_files,
            "scripts": scripts,
            "outputs": [
                f"output/svg/{stem}.svg",
                f"output/pdf/{stem}.pdf",
                f"output/pptx/{stem}_preview.pptx",
                f"output/png/{stem}.png",
            ],
            "generated_by": "deterministic Python scripts",
            "generative_ai_image": False,
            "machine_checked": True,
            "style": "publication (white background, thin axes, no card/poster elements); Fig. 1 is the single schematic figure",
            "author_signoff": "pending before submission",
            "notes": "PPTX is a PNG preview slide, not a fully editable vector source. SVG/PDF plus scripts and source data are canonical.",
        }
    (FIG_ROOT / "figure_provenance.yml").write_text(yaml.safe_dump(provenance, sort_keys=False), encoding="utf-8")
    (OUT / "blend" / "BLENDER_NOT_GENERATED.txt").write_text(
        "Blender is not used for submission figures. The manuscript figure set is 2D publication-style; "
        "Fig. 1 is a flat scientific schematic generated from the archived Python script.\n",
        encoding="utf-8",
    )


def write_readme() -> None:
    (FIG_ROOT / "README_visuals.md").write_text(
        """# Manuscript Figure Package

This folder contains deterministic, editable-first manuscript figure sources.

## Build

```powershell
python src/build_manuscript_figures.py
```

Outputs are written to `output/svg`, `output/pdf`, `output/png`, and `output/pptx`.
SVG/PDF plus scripts and source data are the canonical figure sources; PNG and PPTX files are previews.

## Style contract (publication style)

- Designed at 180 mm (7.05 in) print width; in-figure text 6-8 pt; panel letters bold lowercase.
- White background, thin axes (0.6 pt), outward ticks, no gridlines, frameless legends.
- No in-figure title banners, no rounded-card panels, no drop shadows, no gradients,
  no dashboard/infographic elements.
- Only Fig. 1 is a schematic/mechanism figure; Figs. 2-6 are data figures
  (Fig. 2a, Fig. 4a, and Fig. 5a contain small flat line-schematic panels only).
- Fig. 6 is an R16/R20 candidate closure figure generated for internal review; it is
  not referenced by the current five-figure manuscript unless the text is revised.
- Okabe-Ito colorblind-safe palette: SPT #0072B2, CPT #CC79A7, stable/excluded #009E73,
  critical/caveat #E69F00, failed #D55E00.

## Editing

- Edit SVG files in Inkscape or Illustrator.
- Use the PPTX files only for presentation review; they embed PNG previews and are not fully editable vector artwork.
- Do not redraw quantitative panels in PowerPoint.
- Keep `figure_provenance.yml` synchronized after manual edits.

## Policy Boundary

No generative-AI image model, stock image, web icon, or copyrighted visual asset is used.
All geometry is script-defined and all numerical values come from the manuscript/SI records mirrored in `data/`.
""",
        encoding="utf-8",
    )
    (FIG_ROOT / "Makefile").write_text(
        """env-check:
\tpython -c "import matplotlib, yaml, numpy; print('figure env ok')"

figures:
\tpython src/build_manuscript_figures.py

clean:
\tpython -c "from pathlib import Path; [p.unlink() for d in ['output/svg','output/pdf','output/png','output/pptx'] for p in Path(d).glob('Fig*') if p.is_file()]"
""",
        encoding="utf-8",
    )


def fig6():
    """Fig. 6: pre-registered site-paired Vs critical-band closure (R16/R17).

    Source data: data/fig6_case_level.csv (pooled tier-A out-of-fold
    probabilities from the deterministic closure pipeline archived in
    rounds/R16_vs_critical_band_closure_2026-06-10/). Interval values are
    transcribed from e2_closure_results.json / e2_closure2_results.json.
    """
    case_csv = DATA / "fig6_case_level.csv"
    if not case_csv.exists():
        print("fig6 skipped: data/fig6_case_level.csv missing")
        return
    import pandas as pd
    cl = pd.read_csv(case_csv)

    fig = plt.figure(figsize=(FULL_W, 2.9))
    gs = fig.add_gridspec(1, 3, width_ratios=[1.3, 1.0, 0.72], wspace=0.45)

    # (a) Delta NLL gains across the pre-registered runs
    ax = fig.add_subplot(gs[0, 0])
    panel_letter(ax, "a", dx_pt=-92)
    rows = [
        ("All pooled cases (n=133)", 0.017, -0.042, 0.105, GRAY),
        ("In-band, pooled (primary)", 0.131, 0.079, 0.174, GREEN),
        ("In-band, SPT axis (primary)", 0.120, 0.022, 0.192, GREEN),
        ("In-band, label-agreeing only", 0.272, 0.110, 0.391, BLUE),
        ("In-band, incl. area-level pairs", 0.052, -0.137, 0.160, AMBER),
    ]
    y = np.arange(len(rows))[::-1]
    for yy, (lab, val, lo, hi, col) in zip(y, rows):
        ax.plot([lo, hi], [yy, yy], color=col, lw=1.1, solid_capstyle="butt")
        ax.scatter([val], [yy], color=col, s=14, zorder=3)
    ax.axvline(0, color=DARK, lw=0.6)
    ax.axvline(0.05, color=GRAY, ls="--", lw=0.7)
    ax.text(0.054, 4.45, "+0.05 practical\ngain threshold", fontsize=6.0, color=GRAY, va="top")
    ax.set_yticks(y)
    ax.set_yticklabels([r[0] for r in rows], fontsize=6.5)
    ax.set_xlabel(r"$\Delta$NLL gain from adding the $V_s$ margin")
    ax.set_xlim(-0.16, 0.42)
    ax.set_ylim(-0.5, 4.6)

    # (b) case-level probability shift
    ax = fig.add_subplot(gs[0, 1])
    panel_letter(ax, "b", dx_pt=-30)
    ax.axvspan(0.3, 0.7, color=AMBER, alpha=0.12, lw=0)
    ax.plot([0, 1], [0, 1], color=GRAY, lw=0.7)
    for inst, mk in [("SPT", "o"), ("CPT", "s")]:
        for yval, col in [(1.0, RED), (0.0, GREEN)]:
            sub = cl[(cl["instrument"] == inst) & (cl["y"] == yval)]
            ax.scatter(sub["pA"], sub["pB"], s=11, marker=mk, color=col,
                       edgecolor=DARK, lw=0.25, alpha=0.85)
    ax.set_xlabel("penetration-margin probability")
    ax.set_ylabel("probability with $V_s$ margin")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_xticks([0, 0.3, 0.7, 1.0])
    ax.set_yticks([0, 0.3, 0.7, 1.0])
    handles = [
        Line2D([0], [0], marker="o", color="none", markerfacecolor=RED, markeredgecolor=DARK,
               markeredgewidth=0.25, markersize=3.4, label="liquefied"),
        Line2D([0], [0], marker="o", color="none", markerfacecolor=GREEN, markeredgecolor=DARK,
               markeredgewidth=0.25, markersize=3.4, label="not liquefied"),
        Line2D([0], [0], marker="o", color="none", markerfacecolor="white", markeredgecolor=DARK,
               markeredgewidth=0.5, markersize=3.4, label="SPT pair"),
        Line2D([0], [0], marker="s", color="none", markerfacecolor="white", markeredgecolor=DARK,
               markeredgewidth=0.5, markersize=3.4, label="CPT pair"),
    ]
    ax.legend(handles=handles, loc="upper left", fontsize=5.8, handletextpad=0.3,
              borderaxespad=0.2, labelspacing=0.3)

    # (c) review-burden resolution for in-band cases
    ax = fig.add_subplot(gs[0, 2])
    panel_letter(ax, "c", dx_pt=-34)
    inb = cl[cl["inband_A"]]
    resolved = (inb["pB"] < 0.3) | (inb["pB"] > 0.7)
    correct = ((inb["pB"] > 0.7) & (inb["y"] == 1.0)) | ((inb["pB"] < 0.3) & (inb["y"] == 0.0))
    frac_res = float(resolved.mean())
    frac_corr = float((resolved & correct).mean())
    bars = [("still\nin band", 1 - frac_res, GRAY),
            ("correctly\nresolved", frac_corr, GREEN),
            ("wrongly\nresolved", frac_res - frac_corr, RED)]
    xx = np.arange(3)
    ax.bar(xx, [b[1] for b in bars], 0.55, color=[b[2] for b in bars])
    for x0, b in zip(xx, bars):
        ax.text(x0, b[1] + 0.02, f"{b[1]*100:.0f}%", ha="center", fontsize=6.5, color=DARK)
    ax.set_xticks(xx)
    ax.set_xticklabels([b[0] for b in bars], fontsize=5.2)
    ax.set_ylabel(f"fraction of in-band cases\n(n={len(inb)})")
    ax.set_ylim(0, 0.85)
    save(fig, "Fig6_critical_band_closure", "Fig. 6 Site-paired Vs critical-band closure")


def main() -> None:
    setup()
    write_source_data()
    fig1()
    fig2()
    fig3()
    fig4()
    fig5()
    fig6()
    write_provenance()
    write_readme()
    print("Manuscript figure package generated (publication style)")


if __name__ == "__main__":
    main()
